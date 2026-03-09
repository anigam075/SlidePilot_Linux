from __future__ import annotations

from pathlib import Path

from pptx import Presentation


def _shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    text = shape.text.strip()
    return text


def _extract_notes(slide) -> str:
    notes_slide = getattr(slide, "notes_slide", None)
    if not notes_slide:
        return ""
    notes_tf = getattr(notes_slide, "notes_text_frame", None)
    if not notes_tf:
        return ""
    return (notes_tf.text or "").strip()


def parse_pptx(pptx_path: Path) -> list[dict]:
    presentation = Presentation(str(pptx_path))
    slides: list[dict] = []

    for index, slide in enumerate(presentation.slides, start=1):
        chunks = []
        for shape in slide.shapes:
            text = _shape_text(shape)
            if text:
                chunks.append(text)

        content_text = "\n".join(chunks).strip()
        lines = [line.strip() for line in content_text.splitlines() if line.strip()]
        title = lines[0] if lines else f"Slide {index}"

        slides.append(
            {
                "slide_number": index,
                "title": title,
                "content_text": content_text,
                "notes_text": _extract_notes(slide),
                "image_url": None,
                "script": None,
                "audio_url": None,
            }
        )

    return slides
